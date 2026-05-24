import os
import sys
import json
import tempfile
import io

import fitz          # pip install pymupdf
import pytesseract   # pip install pytesseract
from PIL import Image
from pptx import Presentation  # pip install python-pptx

# ─────────────────────────────────────────────────────────
#  תיקיות
# ─────────────────────────────────────────────────────────
INPUT_FOLDER  = "raw_files"
OUTPUT_FOLDER = os.path.join("data", "json_outputs")


# ─────────────────────────────────────────────────────────
#  מציאת Tesseract (Windows בלבד)
# ─────────────────────────────────────────────────────────
def setup_tesseract():
    if sys.platform != "win32":
        return
    candidates = [
        r"C:\Program Files\Tesseract-OCR\tesseract.exe",
        r"C:\Program Files (x86)\Tesseract-OCR\tesseract.exe",
        r"C:\Users\{}\AppData\Local\Programs\Tesseract-OCR\tesseract.exe".format(
            os.environ.get("USERNAME", "")
        ),
    ]
    for path in candidates:
        if os.path.exists(path):
            pytesseract.pytesseract.tesseract_cmd = path
            return
    raise FileNotFoundError(
        "Tesseract לא נמצא.\n"
        "הורד והתקן מ: https://github.com/UB-Mannheim/tesseract/wiki\n"
        "(בחר 'Add to PATH' בזמן ההתקנה)"
    )

setup_tesseract()


# ─────────────────────────────────────────────────────────
#  ייצוא שקופיות לתמונות דרך PowerPoint COM (Windows)
# ─────────────────────────────────────────────────────────
def export_slides_via_powerpoint(pptx_path, output_dir):
    """
    פותח את ה-PPTX דרך PowerPoint ומייצא כל שקופית כ-PNG.
    מחזיר רשימה ממוינת של נתיבי התמונות.
    """
    import win32com.client  # pip install pywin32

    pptx_path = os.path.abspath(pptx_path)
    output_dir = os.path.abspath(output_dir)

    app = win32com.client.Dispatch("PowerPoint.Application")
    app.Visible = True
    try:
        prs = app.Presentations.Open(pptx_path, ReadOnly=True, Untitled=False, WithWindow=False)
        prs.Export(output_dir, "PNG", 1920, 1080)
        prs.Close()
    finally:
        app.Quit()

    # PowerPoint שומר בשם "Slide1.PNG", "Slide2.PNG" ...
    images = sorted(
        [os.path.join(output_dir, f) for f in os.listdir(output_dir)
         if f.lower().endswith(".png")],
        key=lambda p: int(
            "".join(filter(str.isdigit, os.path.splitext(os.path.basename(p))[0])) or "0"
        )
    )
    return images


# ─────────────────────────────────────────────────────────
#  OCR תמונה → טקסט נקי
# ─────────────────────────────────────────────────────────
FOOTER_NOISE = {
    "algorithms and ds i: sorting",
    "algorithms and ds i",
    "april", "april  2025", "april 2025",
}

def clean_ocr_text(raw_text):
    lines = []
    for line in raw_text.splitlines():
        stripped = line.strip()
        if not stripped:
            continue
        low = stripped.lower()
        if any(low.startswith(n) or low == n for n in FOOTER_NOISE):
            continue
        if stripped.isdigit() and len(stripped) <= 2:
            continue
        lines.append(stripped)
    return "\n".join(lines)


def ocr_image_file(image_path):
    img = Image.open(image_path)
    raw = pytesseract.image_to_string(img, lang="eng")
    return clean_ocr_text(raw)


# ─────────────────────────────────────────────────────────
#  מבנה שקופיות מ-python-pptx (כותרות + הערות)
# ─────────────────────────────────────────────────────────
def get_pptx_structure(file_path):
    prs = Presentation(file_path)
    meta = []
    for slide in prs.slides:
        title = None
        if slide.shapes.title and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text.strip() or None
        notes = None
        if slide.has_notes_slide:
            nf = slide.notes_slide.notes_text_frame
            if nf and nf.text.strip():
                notes = nf.text.strip()
        meta.append({"title": title, "speaker_notes": notes})
    return meta


# ─────────────────────────────────────────────────────────
#  עיבוד PPTX
# ─────────────────────────────────────────────────────────
def process_pptx(file_path):
    slides_meta = get_pptx_structure(file_path)

    with tempfile.TemporaryDirectory() as tmp_dir:
        print(f"    → מייצא שקופיות דרך PowerPoint...")
        image_paths = export_slides_via_powerpoint(file_path, tmp_dir)

        slides_data = []
        for i, meta in enumerate(slides_meta):
            content = ""
            if i < len(image_paths):
                print(f"    → OCR שקופית {i + 1}/{len(slides_meta)}...", end="\r")
                content = ocr_image_file(image_paths[i])

            slides_data.append({
                "page_number": i + 1,
                "title": meta["title"],
                "content": content,
                "speaker_notes": meta["speaker_notes"],
            })

    print()  # newline אחרי ה-\r
    return slides_data


# ─────────────────────────────────────────────────────────
#  עיבוד PDF
# ─────────────────────────────────────────────────────────
def process_pdf(file_path):
    doc = fitz.open(file_path)
    pages_data = []
    for i, page in enumerate(doc):
        # ניסיון ראשון: חילוץ טקסט ישיר
        text = page.get_text("text").strip()
        if len(text) < 50:
            # דף עם תמונות → OCR
            scale = 200 / 72.0
            pix = page.get_pixmap(matrix=fitz.Matrix(scale, scale))
            img = Image.open(io.BytesIO(pix.tobytes("png")))
            text = pytesseract.image_to_string(img, lang="eng").strip()
        if text:
            pages_data.append({"page_number": i + 1, "content": text})
    doc.close()
    return pages_data


# ─────────────────────────────────────────────────────────
#  MAIN
# ─────────────────────────────────────────────────────────
def main():
    if not os.path.exists(INPUT_FOLDER):
        os.makedirs(INPUT_FOLDER)
        print(f"נוצרה תיקייה '{INPUT_FOLDER}'.")
        print(f"שים בה קבצי PDF/PPTX והרץ שוב.")
        return

    os.makedirs(OUTPUT_FOLDER, exist_ok=True)

    files = [f for f in os.listdir(INPUT_FOLDER)
             if os.path.isfile(os.path.join(INPUT_FOLDER, f))]

    if not files:
        print(f"לא נמצאו קבצים ב-'{INPUT_FOLDER}'.")
        return

    print(f"נמצאו {len(files)} קבצים — מתחיל עיבוד...\n")
    processed = skipped = 0

    for file_name in files:
        file_path = os.path.join(INPUT_FOLDER, file_name)
        base_name, ext = os.path.splitext(file_name)
        ext = ext.lower()

        print(f"⏳ מעבד: {file_name}")
        try:
            if ext == ".pptx":
                parsed_data = process_pptx(file_path)
            elif ext == ".pdf":
                parsed_data = process_pdf(file_path)
            else:
                print(f"  ⚠️  סוג קובץ לא נתמך — מדולג.")
                skipped += 1
                continue
        except Exception as e:
            print(f"  ❌ שגיאה: {e}")
            skipped += 1
            continue

        if not parsed_data:
            print(f"  ⚠️  לא נמצא תוכן.")
            skipped += 1
            continue

        out_path = os.path.join(OUTPUT_FOLDER, f"{base_name}.json")
        with open(out_path, "w", encoding="utf-8") as f:
            json.dump(parsed_data, f, ensure_ascii=False, indent=2)

        print(f"  ✅ {base_name}.json ({len(parsed_data)} עמודים)")
        processed += 1

    print(f"\nסיום! עובדו {processed}, דולגו {skipped}.")
    print(f"התוצאות נמצאות ב: {OUTPUT_FOLDER}{os.sep}")


if __name__ == "__main__":
    main()