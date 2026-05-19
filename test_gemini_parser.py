from pptx import Presentation
import json

def extract_text_from_pptx(pptx_path):
    """
    קורא קובץ פאוורפוינט ומחלץ ממנו את כל הטקסט,
    מסודר במבנה נוח לפי מספרי שקופיות.
    """
    print(f"Loading presentation: {pptx_path}...")
    try:
        prs = Presentation(pptx_path)
        slides_data = []
        
        for i, slide in enumerate(prs.slides):
            slide_text = []
            # מעבר על כל האובייקטים בשקופית כדי למצוא תיבות טקסט
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text.strip())
            
            # אם מצאנו טקסט בשקופית, נוסיף אותו לרשימה שלנו
            if slide_text:
                slides_data.append({
                    "slide_number": i + 1,
                    "content": "\n".join(slide_text)
                })
                
        return slides_data
    except Exception as e:
        print(f"Error reading PPTX: {e}")
        return None

# --- הרצת הקוד ---
if __name__ == "__main__":
    file_path = "mavnat1.pptx" 
    
    print("Starting extraction...")
    extracted_slides = extract_text_from_pptx(file_path)
    
    if extracted_slides:
        print(f"Successfully extracted {len(extracted_slides)} slides with text.")
        
        # הדפסת התוכן של השקופית הראשונה כדוגמה למסך
        print("\n--- דוגמה לתוכן השקופית הראשונה ---")
        print(extracted_slides[0]["content"])
        
        # שמירת כל המידע המאורגן לקובץ JSON
        output_file = "parsed_presentation.json"
        with open(output_file, 'w', encoding='utf-8') as f:
            json.dump(extracted_slides, f, ensure_ascii=False, indent=4)
            
        print(f"\n✅ All data formatted and saved successfully to '{output_file}'!")
        print("Now the LLM team can simply load this JSON file.")